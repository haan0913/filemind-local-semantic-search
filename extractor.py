"""
Content Extractor — Multi-format text extraction.

Supports: PDF (PyMuPDF), DOCX, XLSX (basic), PPTX (basic), code/text files.
Handles encoding gracefully with fallback.
"""

import logging
import pathlib
from typing import Callable

try:
    from .protected_secrets import redact_secret_text
except ImportError:
    from protected_secrets import redact_secret_text

logger = logging.getLogger(__name__)

# Extractor registry: ext -> callable(filepath) -> str
EXTRACTORS: dict[str, Callable] = {}


def _sanitize_extracted_text(text: str) -> str:
    """Redact obvious secret-like values before they enter the index."""
    return redact_secret_text(text)


def register_extractor(ext: str):
    """Decorator to register an extractor function for a file extension."""

    def decorator(fn: Callable):
        EXTRACTORS[ext.lower()] = fn
        return fn

    return decorator


# ── Text / Code Files ───────────────────────────────────────────────────────

# Extensions that are read as plain text
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".css",
    ".html",
    ".xml",
    ".json",
    ".toml",
    ".yaml",
    ".yml",
    ".csv",
    ".sql",
    ".sh",
    ".bat",
    ".ps1",
    ".cfg",
    ".ini",
    ".log",
    ".env",
    ".conf",
    ".rst",
    ".tex",
    ".latex",
    ".eml",
}


def _read_text(filepath: str) -> str:
    """Read text file with encoding fallback."""
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(filepath, "r", encoding=enc, errors="replace") as f:
                return f.read()
        except (UnicodeDecodeError, OSError):
            continue
    logger.warning(f"Could not read {filepath} with any encoding")
    return ""


for ext in TEXT_EXTENSIONS:
    EXTRACTORS[ext] = _read_text


# ── PDF (PyMuPDF) ───────────────────────────────────────────────────────────

try:
    import pymupdf

    @register_extractor(".pdf")
    def _extract_pdf(filepath: str) -> str:
        try:
            doc = pymupdf.open(filepath)
            text = "\n\n".join(
                page_text if isinstance(page_text, str) else str(page_text)
                for page_text in (page.get_text() for page in doc)
            )
            doc.close()
            return text
        except Exception as e:
            logger.warning(f"PDF extraction failed for {filepath}: {e}")
            return ""
except ImportError:
    logger.info("PyMuPDF not installed, PDF extraction disabled")


# ── DOCX ────────────────────────────────────────────────────────────────────

try:
    from docx import Document

    @register_extractor(".docx")
    def _extract_docx(filepath: str) -> str:
        try:
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.warning(f"DOCX extraction failed for {filepath}: {e}")
            return ""
except ImportError:
    logger.info("python-docx not installed, DOCX extraction disabled")


# ── XLSX (basic — sheet names + cell values) ────────────────────────────────

try:
    import openpyxl

    @register_extractor(".xlsx")
    def _extract_xlsx(filepath: str) -> str:
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"Sheet: {sheet}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            wb.close()
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"XLSX extraction failed for {filepath}: {e}")
            return ""
except ImportError:
    logger.info("openpyxl not installed, XLSX extraction disabled")


# ── PPTX (basic — slide text and notes) ─────────────────────────────────────

try:
    from pptx import Presentation

    @register_extractor(".pptx")
    def _extract_pptx(filepath: str) -> str:
        try:
            prs = Presentation(filepath)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"Slide {i}")
                for shape in slide.shapes:
                    text_frame = getattr(shape, "text_frame", None)
                    if (
                        getattr(shape, "has_text_frame", False)
                        and text_frame is not None
                    ):
                        for para in text_frame.paragraphs:
                            parts.append(para.text)
                if getattr(slide, "has_notes_slide", False):
                    notes_slide = getattr(slide, "notes_slide", None)
                    notes = getattr(notes_slide, "notes_text_frame", None)
                    if notes is not None and notes.text:
                        parts.append(f"Notes: {notes.text}")
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"PPTX extraction failed for {filepath}: {e}")
            return ""
except ImportError:
    logger.info("python-pptx not installed, PPTX extraction disabled")


# ── EML Emails ──────────────────────────────────────────────────────────────


@register_extractor(".eml")
def _extract_eml(filepath: str) -> str:
    """Read .eml email file and extract text content."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        logger.warning(f"EML extraction failed for {filepath}: {e}")
        return ""


# ── Public API ──────────────────────────────────────────────────────────────


def extract_content(filepath: str, max_size: int = 500_000) -> str:
    """
    Extract text content from a file.

    Args:
        filepath: Path to the file
        max_size: Maximum file size to process (bytes). Returns "" if larger.

    Returns:
        Extracted text content (may be empty on failure)
    """
    path = pathlib.Path(filepath)
    if not path.exists():
        logger.warning(f"File not found: {filepath}")
        return ""

    file_size = path.stat().st_size
    if file_size > max_size:
        logger.info(f"Skipping {filepath} — too large ({file_size} bytes)")
        return ""

    name = path.name.lower()
    ext = ".env" if name == ".env" or name.startswith(".env.") else path.suffix.lower()
    extractor = EXTRACTORS.get(ext)

    if extractor:
        try:
            content = extractor(str(path))
            return _sanitize_extracted_text(content)
        except Exception as e:
            logger.error(f"Extraction error for {filepath}: {e}")
            return ""
    else:
        logger.debug(f"No extractor for extension: {ext}")
        return ""


def get_supported_extensions() -> list[str]:
    """Return list of supported file extensions."""
    return sorted(EXTRACTORS.keys())
